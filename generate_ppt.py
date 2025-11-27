from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 读取数据
df = pd.read_csv('runs_summary_statistics.csv')

# 标题页
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "REINVENT4 Multi-Run Analysis"
subtitle.text = "DENV NS2B-NS3 Protease Inhibitor Discovery"

# 摘要页
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Summary Statistics"

# 添加表格
rows, cols = df.shape
table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(2), 
                                Inches(9), Inches(4)).table

# 填充表格
for j, col in enumerate(df.columns):
    table.cell(0, j).text = col
    
for i in range(rows):
    for j in range(cols):
        table.cell(i+1, j).text = str(df.iloc[i, j])

prs.save('runs_analysis.pptx')
print("✅ PowerPoint已生成: runs_analysis.pptx")
